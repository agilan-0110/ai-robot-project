from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from sentence_transformers import SentenceTransformer
from docx import Document
from pypdf import PdfReader, PdfWriter
from markitdown import MarkItDown
import numpy as np
import os
import re
import time
import base64
import tempfile
import requests
import unicodedata
import io
from PIL import Image

# ——— Jetson / Device Configuration ———
# By default, use CPU for sentence-transformers on Jetson to keep all 8GB
# unified memory completely available for Whisper and YOLOv8-pose.
RAG_DEVICE = os.environ.get("RAG_EMBEDDING_DEVICE", "cpu")

# ——— Unicode and text normalization replacements ———
UNICODE_REPLACEMENTS = {
    '\u2018': "'", '\u2019': "'",   # smart single quotes
    '\u201c': '"', '\u201d': '"',   # smart double quotes
    '\u2013': '-', '\u2014': ' - ', # en/em dashes
    '\u2011': '-',                  # non-breaking hyphen
    '\u00a0': ' ',                  # non-breaking space
    '\u2026': '...',                # horizontal ellipsis
    '\u2022': ' ', '\u25ba': ' ',   # bullets & play arrows
}

def normalize_text(text):
    """
    Applies NFKC normalization and replaces problematic Unicode punctuation
    (smart quotes, non-breaking hyphens, non-breaking spaces) with standard ASCII equivalents.
    Ensures safe handling by downstream LLM prompts, TTS, and console logs across platforms.
    """
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    for bad, good in UNICODE_REPLACEMENTS.items():
        text = text.replace(bad, good)
    # Strip unprintable control characters (keep \n, \r, \t)
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    return text.strip()


# ——— Image captioning via Gemini vision model ———
# Uses a separate GEMINI_API_KEY env var (get one free at aistudio.google.com).
# CONFIRMED WORKING (Aug 2026) via manual testing: gemini-3.5-flash-lite —
# ~1.8s/image, good caption quality. The non-lite gemini-3.5-flash also
# works but is ~10x slower (~17s/image) for similar quality — not worth
# it for a per-image upload-time cost. If this model gets deprecated
# later, run test_gemini_vision.py --list-models to find the current one.
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
GEMINI_VISION_MODEL = "gemini-3.5-flash-lite"
GEMINI_BASE_URL = "https://generativelanguage.googleapis.com/v1beta"
IMAGE_CAPTION_PROMPT = (
    "Describe this lecture slide image or diagram in 1-2 sentences. "
    "Focus on the educational content: labels, steps, relationships, or "
    "data shown. If it's purely decorative (logo, background), say so briefly."
)

# ═══════════════════════════════════════════════════════════════════
# ——— Irrelevant content filter — EDIT THIS LIST ANYTIME ———
# Any line matching one of these phrases (AND being short — see
# MAX_WORDS_FOR_IRRELEVANT_LINE below) gets dropped from extraction:
# title-slide clutter like student name, register no, department,
# guide name, etc. This is NOT code logic — just add/remove plain
# phrases here as you test on real decks and notice something slip
# through or get wrongly filtered. No need to touch anything else.
# ═══════════════════════════════════════════════════════════════════
IRRELEVANT_PHRASES = [
    "submitted by", "presented by", "prepared by", "compiled by",
    "guided by", "under the guidance of", "guide name", "guide:",
    "department of", "dept of", "dept.",
    "register no", "register number", "reg no", "reg. no", "reg number",
    "roll no", "roll number",
    "academic year", "batch:", "section:", "class:",
    "student name", "name:", "name ·", "name -",
    "college name", "university name", "affiliated to",
    "under the supervision of", "hod", "head of department",
]

# A line is only filtered if BOTH signals agree: it matches a phrase
# above AND it's short. A long real sentence that happens to loosely
# contain one of these words (rare, but possible) is left alone —
# only short, clearly-administrative lines get dropped.
MAX_WORDS_FOR_IRRELEVANT_LINE = 12

# A line containing a long run of digits (6+ in a row) is almost always
# a register/roll number — catches those even with unpredictable
# label wording (e.g. "SUBASH M · 3RD YR AI&DS / A · 113224072105"),
# without needing to guess every possible phrasing in the list above.
_LONG_DIGIT_RUN_PATTERN = re.compile(r'\d{6,}')


def _is_irrelevant_line(line):
    """
    Returns True if this line looks like title-slide clutter
    (name/department/register no/etc.) rather than real lecture
    content, and should be dropped from extraction.
    """
    stripped = line.strip()
    if not stripped:
        return False

    word_count = len(stripped.split())
    if word_count > MAX_WORDS_FOR_IRRELEVANT_LINE:
        return False  # long lines are never filtered — too risky to be wrong

    lower = stripped.lower()
    if any(phrase in lower for phrase in IRRELEVANT_PHRASES):
        return True
    if _LONG_DIGIT_RUN_PATTERN.search(stripped):
        return True

    return False

# ——— Global state (loaded once, reused across requests) ———
_model = None
_slide_texts = []   # each item also carries "source_file"
_vectors = None
_loaded_files = []  # list of file paths currently loaded (supports multiple)
_conversation_history = []  # list of {"question", "answer", "chunks"} for this class session
_current_lecture_slide = 0
_max_slide_reached = 0

_FOLLOWUP_PHRASES = [
    "explain again", "explain that again", "explain one more time",
    "one more time", "say that again", "say again", "repeat that", "repeat it",
    "please repeat", "please repeat that",
    "didn't understand", "didnt understand", "didn't get it", "didnt get it",
    "not clear", "not clear sir", "not understood", "unclear",
    "i didn't get it", "i didnt get it", "still confused", "i m confused",
    "i am confused", "im confused", "confused",
    "come again", "come again sir", "once more", "again please", "again sir",
    "explain more", "can you explain more", "i'm lost", "im lost",
]

# ——— Explicit slide-number detection ———
# Semantic (embedding) search has no real concept of "slide 5" — the
# embedding of the literal words "explain slide 5" isn't reliably close
# to whatever content actually lives on slide 5. So a direct question
# like "explain slide 5" must be resolved by exact slide_number lookup,
# not similarity search, or it ends up answering an unrelated slide.
_WORD_TO_NUM = {
    "one": 1, "two": 2, "three": 3, "four": 4, "five": 5, "six": 6,
    "seven": 7, "eight": 8, "nine": 9, "ten": 10, "eleven": 11,
    "twelve": 12, "thirteen": 13, "fourteen": 14, "fifteen": 15,
    "sixteen": 16, "seventeen": 17, "eighteen": 18, "nineteen": 19,
    "twenty": 20, "thirty": 30, "forty": 40, "fifty": 50,
}
_SLIDE_NUMBER_DIGIT_PATTERN = re.compile(r'slide\s*(?:number\s*|#\s*)?(\d+)', re.IGNORECASE)
_SLIDE_NUMBER_WORD_PATTERN = re.compile(
    r'slide\s*(?:number\s*)?(' + "|".join(_WORD_TO_NUM.keys()) + r')', re.IGNORECASE
)


def extract_requested_slide_number(text):
    """
    Returns the slide number explicitly mentioned in the text (as an int),
    or None if no slide reference is found. Handles both digit form
    ("slide 5") and spoken/word form ("slide five") — Whisper frequently
    transcribes small spoken numbers as words rather than digits.
    """
    match = _SLIDE_NUMBER_DIGIT_PATTERN.search(text)
    if match:
        return int(match.group(1))

    match = _SLIDE_NUMBER_WORD_PATTERN.search(text.lower())
    if match:
        return _WORD_TO_NUM.get(match.group(1))

    return None


_markitdown = None


def _get_model():
    global _model
    if _model is None:
        print(f"[RAG] Loading embedding model on {RAG_DEVICE}...")
        _model = SentenceTransformer("all-MiniLM-L6-v2", device=RAG_DEVICE)
    return _model


def _get_markitdown():
    global _markitdown
    if _markitdown is None:
        _markitdown = MarkItDown()
    return _markitdown


def _is_decorative_image(image_bytes):
    """
    Returns True if the image appears to be a small icon, bullet, or
    decorative graphic rather than an educational diagram, chart, or photo.
    Filtering these out:
      1. Accelerates upload ingestion by 4-5x.
      2. Prevents Gemini 15 RPM rate limits on decks with lots of icons.
      3. Keeps the RAG knowledge base focused purely on educational content.
    """
    if not image_bytes or len(image_bytes) < 2048:  # under 2KB is virtually always a tiny icon/bullet
        return True
    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            w, h = img.size
            # Skip tiny dimensions (bullets, status dots, small UI icons)
            if w < 80 or h < 80:
                return True
            # Skip small square-ish icons (< 25,000 sq px and < 15KB)
            if (w * h < 25000) and len(image_bytes) < 15000:
                return True
    except Exception:
        pass
    return False


def _caption_image_bytes(image_bytes, mime_type="image/png"):
    """
    Sends raw image bytes to Gemini and returns a short caption.
    Fails soft (returns "") if no API key is set or the call errors
    after retries — a missing caption should never block a file upload.

    Retries on 503 (Gemini high demand) and 429 (rate limit exceeded)
    with backoff before giving up.
    """
    if not GEMINI_API_KEY:
        return ""

    b64 = base64.b64encode(image_bytes).decode("utf-8")
    url = f"{GEMINI_BASE_URL}/models/{GEMINI_VISION_MODEL}:generateContent?key={GEMINI_API_KEY}"
    payload = {
        "contents": [{
            "parts": [
                {"text": IMAGE_CAPTION_PROMPT},
                {"inline_data": {"mime_type": mime_type, "data": b64}},
            ]
        }]
    }

    max_attempts = 3
    for attempt in range(max_attempts):
        try:
            resp = requests.post(url, json=payload, timeout=20)
            if resp.status_code == 503 and attempt < max_attempts - 1:
                print(f"[RAG] Gemini busy (503), retrying in 2s... (attempt {attempt + 1}/{max_attempts})")
                time.sleep(2)
                continue
            if resp.status_code == 429 and attempt < max_attempts - 1:
                retry_after = resp.headers.get("Retry-After")
                wait_sec = float(retry_after) if retry_after else 4.0 * (attempt + 1)
                print(f"[RAG] Gemini rate limit (429), retrying in {wait_sec:.1f}s... (attempt {attempt + 1}/{max_attempts})")
                time.sleep(wait_sec)
                continue
            resp.raise_for_status()
            data = resp.json()
            return data["candidates"][0]["content"]["parts"][0]["text"].strip()
        except Exception as e:
            if attempt == max_attempts - 1:
                print(f"[RAG] Image captioning failed after {max_attempts} attempts (skipping image): {e}")
                return ""
            print(f"[RAG] Image captioning attempt {attempt + 1} failed ({e}), retrying in 1s...")
            time.sleep(1)
            continue
    return ""


def _caption_pptx_slide_images(slide):
    """Returns a list of '[Image: ...]' caption strings for every real
    embedded picture on this slide, skipping decorative icons and SmartArt.
    """
    captions = []
    skipped_non_embedded = 0
    skipped_icons = 0
    picture_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    real_pictures = []
    for s in picture_shapes:
        try:
            image = s.image
            if _is_decorative_image(image.blob):
                skipped_icons += 1
                continue
            real_pictures.append(image)
        except ValueError:
            skipped_non_embedded += 1
            continue

    for idx, image in enumerate(real_pictures, start=1):
        print(f"[RAG]   -> Captioning image {idx}/{len(real_pictures)} via Gemini...")
        try:
            caption = _caption_image_bytes(image.blob, image.content_type or "image/png")
        except Exception as e:
            print(f"[RAG] Could not caption pptx image: {e}")
            caption = ""
        if caption:
            print(f"[RAG]      Caption: {caption[:70]}...")
            captions.append(f"[Image: {caption}]")

    if skipped_icons:
        print(f"[RAG]   -> Skipped {skipped_icons} small/decorative icon(s) (kept RAG clean)")
    if skipped_non_embedded:
        print(f"[RAG]   -> Skipped {skipped_non_embedded} SmartArt/diagram graphic(s) "
              f"(not real photos, their text is already captured separately)")

    return captions


def _caption_pdf_page_images(page):
    """Returns a list of '[Image: ...]' caption strings for every real
    educational image embedded on this PDF page."""
    captions = []
    skipped_icons = 0
    real_images = []
    for img in page.images:
        try:
            if _is_decorative_image(img.data):
                skipped_icons += 1
                continue
            real_images.append(img)
        except Exception:
            continue

    for idx, img in enumerate(real_images, start=1):
        try:
            caption = _caption_image_bytes(img.data, "image/png")
            if caption:
                captions.append(f"[Image: {caption}]")
        except Exception as e:
            print(f"[RAG] Could not read pdf image: {e}")

    if skipped_icons:
        print(f"[RAG]   -> Skipped {skipped_icons} small/decorative icon(s) on PDF page")
    return captions


def _caption_docx_images(path):
    """
    Captions educational images in DOCX, skipping decorative icons.
    """
    doc = Document(path)
    captions = []
    skipped_icons = 0
    real_images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                if _is_decorative_image(blob):
                    skipped_icons += 1
                    continue
                real_images.append((blob, rel.target_part.content_type))
            except Exception as e:
                print(f"[RAG] Could not read docx image: {e}")

    for blob, content_type in real_images:
        try:
            caption = _caption_image_bytes(blob, content_type)
            if caption:
                captions.append(f"[Image: {caption}]")
        except Exception as e:
            print(f"[RAG] Could not read docx image: {e}")

    if skipped_icons:
        print(f"[RAG]   -> Skipped {skipped_icons} small/decorative icon(s) in DOCX")
    return captions


def _strip_markitdown_comments(md_text):
    """Removes MarkItDown's own '<!-- Slide number: N -->' comments AND
    its bare '![filename](Picture2.jpg)' image placeholders — neither
    carries real information, and we don't caption images (no vision
    model wired in currently), so any embedded images are simply
    excluded from the extracted text."""
    lines = []
    for l in md_text.split("\n"):
        stripped = l.strip()
        if stripped.startswith("<!--") and stripped.endswith("-->"):
            continue
        if stripped.startswith("!["):
            continue
        lines.append(l)
    return "\n".join(lines).strip()


GENERIC_HEADINGS = {
    "notes", "notes:", "speaker notes", "speaker notes:", "slide notes",
    "untitled", "untitled slide", "agenda", "overview", "table of contents",
}
_GENERIC_HEADING_REGEX = re.compile(r'^(?:slide\s*(?:number\s*|#\s*)?\d+|page\s*\d+)$', re.IGNORECASE)


def _is_generic_heading(candidate):
    """Returns True if the heading is a generic placeholder like 'Notes:', '### Notes:', or 'Slide 1'."""
    clean = candidate.lstrip("#").strip().lower()
    if clean in GENERIC_HEADINGS:
        return True
    if _GENERIC_HEADING_REGEX.match(clean):
        return True
    return False


def _markdown_to_heading_and_points(md_text):
    """
    Splits a single chunk's markdown into (heading, points):
    - heading = first non-generic '#' line found, if any
    - points  = every other non-empty, non-irrelevant line
    All lines are normalized (NFKC and clean punctuation).
    """
    heading = ""
    points = []
    for raw_line in md_text.split("\n"):
        line = normalize_text(raw_line)
        if not line:
            continue
        if line.startswith("<!--") and line.endswith("-->"):
            continue
        if line.startswith("!["):
            continue
        if line.startswith("#") and not heading:
            candidate_heading = line.lstrip("#").strip()
            if not _is_irrelevant_line(candidate_heading) and not _is_generic_heading(candidate_heading):
                heading = candidate_heading
        else:
            if _is_generic_heading(line):
                continue
            clean_line = line.lstrip("#").strip() if line.startswith("#") else line
            if clean_line and not _is_irrelevant_line(clean_line) and not _is_generic_heading(clean_line):
                points.append(clean_line)
    return heading, points


def _build_chunk_text(heading, points):
    """
    Reconstructs the clean 'text' field FROM the already-filtered, normalized
    heading/points so irrelevant clutter and generic markers are completely removed.
    """
    clean_heading = normalize_text(heading)
    clean_points = [normalize_text(p) for p in points if normalize_text(p)]
    lines = ([f"# {clean_heading}"] if clean_heading else []) + clean_points
    return "\n".join(lines).strip()


def _extract_from_pptx(path):
    """
    Converts each slide individually via MarkItDown with per-slide fault tolerance.
    """
    converter = _get_markitdown()
    prs = Presentation(path)
    num_slides = len(prs.slides)
    chunks = []

    with tempfile.TemporaryDirectory() as tmp_dir:
        for i in range(num_slides):
            print(f"[RAG] Processing slide {i + 1}/{num_slides}...")
            try:
                single = Presentation(path)
                slide_id_list = single.slides._sldIdLst
                all_ids = list(slide_id_list)
                for j, sld in enumerate(all_ids):
                    if j != i:
                        slide_id_list.remove(sld)
                slide_path = os.path.join(tmp_dir, f"_slide_{i}.pptx")
                single.save(slide_path)

                result = converter.convert(slide_path)
                md_text = _strip_markitdown_comments(result.text_content or "")

                # Caption any pictures on the real (original) slide object
                image_captions = _caption_pptx_slide_images(prs.slides[i])
                if image_captions:
                    md_text = (md_text + "\n" + "\n".join(image_captions)).strip()

                if not md_text:
                    continue

                heading, points = _markdown_to_heading_and_points(md_text)
                clean_text = _build_chunk_text(heading, points)
                if not clean_text:
                    continue  # entire slide was filler

                chunks.append({
                    "slide_number": i + 1,
                    "text": clean_text,
                    "heading": heading,
                    "points": points,
                })
            except Exception as e:
                print(f"[RAG] Warning: Slide {i + 1} could not be extracted ({e}) — skipping.")
                continue
    return chunks


def _extract_from_pdf(path):
    """
    Converts each page individually via MarkItDown with per-page fault tolerance.
    """
    converter = _get_markitdown()
    reader = PdfReader(path)
    num_pages = len(reader.pages)
    chunks = []

    with tempfile.TemporaryDirectory() as tmp_dir:
        for i in range(num_pages):
            print(f"[RAG] Processing page {i + 1}/{num_pages}...")
            try:
                writer = PdfWriter()
                writer.add_page(reader.pages[i])
                page_path = os.path.join(tmp_dir, f"_page_{i}.pdf")
                with open(page_path, "wb") as f:
                    writer.write(f)

                result = converter.convert(page_path)
                md_text = _strip_markitdown_comments(result.text_content or "")

                image_captions = _caption_pdf_page_images(reader.pages[i])
                if image_captions:
                    md_text = (md_text + "\n" + "\n".join(image_captions)).strip()

                if not md_text:
                    continue

                heading, points = _markdown_to_heading_and_points(md_text)
                clean_text = _build_chunk_text(heading, points)
                if not clean_text:
                    continue

                chunks.append({
                    "slide_number": i + 1,
                    "text": clean_text,
                    "heading": heading,
                    "points": points,
                })
            except Exception as e:
                print(f"[RAG] Warning: Page {i + 1} could not be extracted ({e}) — skipping.")
                continue
    return chunks


def _extract_from_docx(path):
    """
    Converts DOCX via MarkItDown and groups by heading with fault tolerance.
    """
    converter = _get_markitdown()
    try:
        result = converter.convert(path)
        md_text = _strip_markitdown_comments(result.text_content or "")
    except Exception as e:
        print(f"[RAG] Warning: Could not convert docx file ({e})")
        return []

    if not md_text:
        return []

    chunks = []
    current_lines = []
    chunk_number = 1

    def flush():
        nonlocal current_lines, chunk_number
        text_block = "\n".join(l.strip() for l in current_lines if l.strip())
        if text_block:
            heading, points = _markdown_to_heading_and_points(text_block)
            clean_text = _build_chunk_text(heading, points)
            if clean_text:
                chunks.append({
                    "slide_number": chunk_number,
                    "text": clean_text,
                    "heading": heading,
                    "points": points,
                })
                chunk_number += 1
        current_lines = []

    for raw_line in md_text.split("\n"):
        line = raw_line.strip()
        if line.startswith("#") and current_lines:
            flush()
        current_lines.append(line)
    flush()

    image_captions = _caption_docx_images(path)
    if image_captions:
        text_block = "Images in this document:\n" + "\n".join(image_captions)
        chunks.append({
            "slide_number": chunk_number,
            "text": text_block,
            "heading": "Document Images",
            "points": image_captions,
        })

    return chunks


def _extract_text(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        return _extract_from_pptx(path)
    elif ext == ".pdf":
        return _extract_from_pdf(path)
    elif ext == ".docx":
        return _extract_from_docx(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def load_lecture(file_path, append=False):
    """
    Public interface: processes a lecture file (pptx/pdf/docx) into
    embeddings, held in memory for retrieval.

    append=False: replaces whatever was loaded before.
    append=True: adds this file's content alongside whatever is
        already loaded (multiple files in one class session).
    """
    global _slide_texts, _vectors, _loaded_files, _conversation_history

    print(f"[RAG] Loading lecture: {file_path} (append={append})")
    chunks = _extract_text(file_path)

    if not chunks:
        raise ValueError("No text could be extracted from this file.")

    source_name = os.path.basename(file_path)
    for c in chunks:
        c["source_file"] = source_name

    model = _get_model()
    texts = [c["text"] for c in chunks]
    new_vectors = model.encode(texts, show_progress_bar=False)

    if append and _vectors is not None:
        _slide_texts = _slide_texts + chunks
        _vectors = np.vstack([_vectors, new_vectors])
        _loaded_files.append(file_path)
    else:
        _slide_texts = chunks
        _vectors = new_vectors
        _loaded_files = [file_path]
        _conversation_history = []  # fresh class, fresh memory (only on non-append load)

    print(f"[RAG] Loaded {len(chunks)} chunks from {source_name} "
          f"(total chunks now: {len(_slide_texts)})")
    return len(chunks)


def remove_file(file_path):
    """
    Public interface: removes a single loaded file's chunks/vectors
    from memory without touching any other loaded files. Does NOT
    touch conversation history (a past answer stays valid even if
    the source file is later removed).

    Returns True if the file was found and removed, False otherwise.
    """
    global _slide_texts, _vectors, _loaded_files

    if file_path not in _loaded_files:
        return False

    source_name = os.path.basename(file_path)

    # build a mask of which chunks/vectors to KEEP (i.e. NOT from this file)
    keep_indices = [
        i for i, c in enumerate(_slide_texts)
        if c.get("source_file") != source_name
    ]

    if keep_indices:
        _slide_texts = [_slide_texts[i] for i in keep_indices]
        _vectors = _vectors[keep_indices]
    else:
        _slide_texts = []
        _vectors = None

    _loaded_files = [f for f in _loaded_files if f != file_path]

    print(f"[RAG] Removed file: {file_path} (remaining chunks: {len(_slide_texts)})")
    return True


def retrieve_relevant_slides(question, top_k=3):
    """
    Public interface: given a student's question, returns the top_k
    most relevant chunks across all currently loaded files, as a list
    of {"slide_number": int, "text": str, "heading": str, "score": float, "source_file": str}.
    When _max_slide_reached > 0, filters out any candidate slide where slide_number > _max_slide_reached.
    Returns an empty list if no lecture is loaded yet or question is empty.
    """
    if _vectors is None or not _slide_texts or top_k <= 0:
        return []

    clean_q = normalize_text(question)
    if not clean_q:
        return []

    # Filter candidate indices to only those slides covered so far (if lecture tracking is active)
    valid_indices = [
        i for i, c in enumerate(_slide_texts)
        if _max_slide_reached <= 0 or c.get("slide_number", 0) <= _max_slide_reached
    ]
    if not valid_indices:
        return []

    model = _get_model()
    question_vector = model.encode([clean_q], show_progress_bar=False)[0]

    q_mag = np.linalg.norm(question_vector)
    if q_mag < 1e-9:
        return []
    q_norm = question_vector / (q_mag + 1e-9)

    candidate_vectors = _vectors[valid_indices]
    v_mags = np.linalg.norm(candidate_vectors, axis=1, keepdims=True)
    v_mags[v_mags < 1e-9] = 1e-9  # Avoid division by zero
    v_norm = candidate_vectors / v_mags

    similarities = np.dot(v_norm, q_norm)
    similarities = np.nan_to_num(similarities, nan=-1.0)

    effective_k = min(top_k, len(valid_indices))
    top_sub_indices = np.argsort(similarities)[::-1][:effective_k]

    results = []
    for sub_idx in top_sub_indices:
        orig_idx = valid_indices[sub_idx]
        results.append({
            "slide_number": _slide_texts[orig_idx]["slide_number"],
            "text": _slide_texts[orig_idx]["text"],
            "heading": _slide_texts[orig_idx].get("heading", ""),
            "score": float(similarities[sub_idx]),
            "source_file": _slide_texts[orig_idx].get("source_file", "unknown"),
        })
    return results


def _is_followup(question):
    q = question.lower().strip()

    if any(phrase in q for phrase in _FOLLOWUP_PHRASES):
        return True

    word_count = len(q.split())
    if word_count <= 3 and _conversation_history:
        quick_check = retrieve_relevant_slides(question, top_k=1)
        if quick_check and quick_check[0]["score"] < 0.35:
            return True

    return False


def _build_lecture_text(chunks):
    """
    Joins retrieved chunks into a single pre-formatted block, ready to
    drop straight into an LLM prompt — so callers (e.g. Nithesh's
    get_llm_answer()) don't need to duplicate this formatting logic.
    """
    return "\n\n".join(
        f"[{c['source_file']}, slide {c['slide_number']}]: {c['text']}"
        for c in chunks
    )


def get_context_for_query(question, top_k=3):
    """
    Public interface: the 'smart' entry point the orchestrator should
    call instead of retrieve_relevant_slides() directly.
    """
    # Explicit "explain slide N" / "slide five" references must be
    # resolved by exact lookup, not semantic search — see
    # extract_requested_slide_number()'s docstring for why.
    requested_slide = extract_requested_slide_number(question)
    if requested_slide is not None:
        matching_chunks = [c for c in _slide_texts if c["slide_number"] == requested_slide]
        if matching_chunks:
            # If the slide exists in the deck but hasn't been taught yet in the live lecture
            if _max_slide_reached > 0 and requested_slide > _max_slide_reached:
                print(f"[RAG] Slide {requested_slide} requested but has not been taught yet (max reached: {_max_slide_reached}).")
                return {
                    "is_followup": False,
                    "not_covered_yet": True,
                    "requested_slide": requested_slide,
                    "chunks": [],
                    "previous_question": None,
                    "previous_answer": None,
                    "lecture_text": f"Slide {requested_slide} has not been covered yet in today's lecture.",
                }

            # If multiple files loaded, check if question mentions a specific file name;
            # otherwise prioritize the most recently loaded deck
            chosen = matching_chunks[-1]
            q_lower = question.lower()
            for c in matching_chunks:
                src = c.get("source_file", "").lower()
                base_src, _ = os.path.splitext(src)
                tokens = [t for t in re.split(r'[-_ ]+', base_src) if len(t) > 2]
                if src in q_lower or any(t in q_lower for t in tokens):
                    chosen = c
                    break

            print(f"[RAG] Explicit slide reference detected -> slide {requested_slide} ({chosen.get('source_file')}) (exact match, no semantic search)")
            chunks = [{
                "slide_number": chosen["slide_number"],
                "text": chosen["text"],
                "heading": chosen.get("heading", ""),
                "score": 1.0,
                "source_file": chosen.get("source_file", "unknown"),
            }]
            return {
                "is_followup": False,
                "not_covered_yet": False,
                "requested_slide": requested_slide,
                "chunks": chunks,
                "previous_question": None,
                "previous_answer": None,
                "lecture_text": _build_lecture_text(chunks),
            }
        else:
            print(f"[RAG] Slide {requested_slide} was requested but doesn't exist in loaded lecture(s) "
                  f"({len(_slide_texts)} slides loaded) — falling back to semantic search.")

    if _is_followup(question) and _conversation_history:
        last = _conversation_history[-1]
        return {
            "is_followup": True,
            "not_covered_yet": False,
            "requested_slide": None,
            "chunks": last["chunks"],
            "previous_question": last["question"],
            "previous_answer": last["answer"],
            "lecture_text": _build_lecture_text(last["chunks"]),
        }

    chunks = retrieve_relevant_slides(question, top_k=top_k)
    return {
        "is_followup": False,
        "not_covered_yet": False,
        "requested_slide": None,
        "chunks": chunks,
        "previous_question": None,
        "previous_answer": None,
        "lecture_text": _build_lecture_text(chunks),
    }


# Alias — the integration handoff doc refers to this function as
# retrieve_context(); both names now work so either the orchestrator
# or this module can use its preferred name without breaking the other.
retrieve_context = get_context_for_query


def add_to_history(question, answer, chunks):
    """
    Public interface: call this after Groq generates an answer, so
    future follow-ups ("explain again") can reuse this turn.
    """
    _conversation_history.append({
        "question": question,
        "answer": answer,
        "chunks": chunks,
    })


def clear_lecture():
    """Public interface: wipes ALL currently loaded lecture files AND conversation history."""
    global _slide_texts, _vectors, _loaded_files, _conversation_history, _current_lecture_slide, _max_slide_reached
    _slide_texts = []
    _vectors = None
    _loaded_files = []
    _conversation_history = []
    _current_lecture_slide = 0
    _max_slide_reached = 0
    print("[RAG] Lecture(s) and conversation history cleared from memory.")


def has_lecture_loaded():
    return _vectors is not None


def get_loaded_files():
    """Public interface: returns list of currently loaded file paths."""
    return list(_loaded_files)


def get_ordered_chunks():
    """
    Public interface: returns all currently loaded chunks in their
    original slide/page order, for the automatic 'teach the class'
    walkthrough (as opposed to retrieve_relevant_slides(), which
    returns only the top-k matches for a specific question).
    """
    return list(_slide_texts)


def set_lecture_progress(current_slide, max_slide=None):
    """
    Public interface: tracks the main lecture's current position and max slide reached.
    _current_lecture_slide only changes on 'next', never on doubt-detours.
    _max_slide_reached only ever increases, never decreases.
    """
    global _current_lecture_slide, _max_slide_reached
    _current_lecture_slide = current_slide
    if max_slide is not None:
        _max_slide_reached = max(_max_slide_reached, max_slide)
    else:
        _max_slide_reached = max(_max_slide_reached, current_slide)


def get_lecture_progress():
    """Public interface: returns (current_lecture_slide, max_slide_reached)."""
    return _current_lecture_slide, _max_slide_reached


def set_max_slide_reached(max_slide):
    """Public interface: sets max_slide_reached directly."""
    global _max_slide_reached
    _max_slide_reached = max_slide


def get_max_slide_reached():
    """Public interface: gets max_slide_reached."""
    return _max_slide_reached