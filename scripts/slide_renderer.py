"""
Slide Renderer Module
---------------------
Converts original presentation decks (.pptx, .pdf) into high-resolution visual
slide images (PNG) so the classroom projector displays the exact original slides
with full graphics, fonts, colors, and layout.
"""

import os
import sys
import shutil
import pypdfium2
from PIL import Image

CACHE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static", "slide_cache")
os.makedirs(CACHE_DIR, exist_ok=True)


def get_deck_cache_dir(file_path):
    """Returns the directory path for storing rendered slide images of a file."""
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    safe_name = "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in base_name)
    deck_dir = os.path.join(CACHE_DIR, safe_name)
    os.makedirs(deck_dir, exist_ok=True)
    return deck_dir


def render_pdf_to_images(pdf_path, out_dir):
    """Renders all pages of a PDF into high-res PNG images."""
    pdf = pypdfium2.PdfDocument(pdf_path)
    image_paths = []
    for i, page in enumerate(pdf, start=1):
        # Render at 2x scale for 1080p/4K crispness
        pil_img = page.render(scale=2.0).to_pil()
        out_path = os.path.join(out_dir, f"slide_{i}.png")
        pil_img.save(out_path, "PNG")
        image_paths.append(out_path)
    return image_paths


def render_pptx_to_images(pptx_path, out_dir):
    """Renders all slides of a PPTX into 1920x1080 PNG images."""
    abs_pptx = os.path.abspath(pptx_path)
    abs_out = os.path.abspath(out_dir)
    image_paths = []

    # 1. On Windows with Office PowerPoint installed: Use COM automation
    if sys.platform == "win32":
        try:
            import comtypes
            import comtypes.client
            comtypes.CoInitialize()
            try:
                powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                pres = powerpoint.Presentations.Open(abs_pptx, True, False, False)
                for i, slide in enumerate(pres.Slides, start=1):
                    out_path = os.path.join(abs_out, f"slide_{i}.png")
                    slide.Export(out_path, "PNG", 1920, 1080)
                    image_paths.append(out_path)
                pres.Close()
                powerpoint.Quit()
            finally:
                comtypes.CoUninitialize()
            return image_paths
        except Exception as e:
            print(f"[RENDERER] Windows COM export error ({e}), trying fallback...")

    # 2. LibreOffice fallback (for Linux / Jetson / systems with LibreOffice)
    try:
        import subprocess
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", abs_out, abs_pptx],
            check=True,
            timeout=30
        )
        base_name = os.path.splitext(os.path.basename(abs_pptx))[0]
        temp_pdf = os.path.join(abs_out, f"{base_name}.pdf")
        if os.path.exists(temp_pdf):
            image_paths = render_pdf_to_images(temp_pdf, abs_out)
            try:
                os.remove(temp_pdf)
            except Exception:
                pass
            return image_paths
    except Exception as e:
        print(f"[RENDERER] LibreOffice conversion error: {e}")

    # 3. Native Python PIL / python-pptx fallback (guaranteed on any system)
    return render_pptx_fallback_pil(abs_pptx, abs_out)


def render_pptx_fallback_pil(pptx_path, out_dir):
    """Fallback visual slide generator using python-pptx & Pillow."""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw
        prs = Presentation(pptx_path)
        image_paths = []
        os.makedirs(out_dir, exist_ok=True)
        for i, slide in enumerate(prs.slides, start=1):
            img = Image.new("RGB", (1920, 1080), color=(15, 23, 42))
            draw = ImageDraw.Draw(img)

            title = ""
            body_lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            if not title:
                                title = txt
                            else:
                                body_lines.append(txt)

            draw.rectangle([(0, 0), (1920, 140)], fill=(30, 41, 59))
            draw.text((80, 45), f"Slide {i}: {title or 'Lecture Slide'}", fill=(248, 250, 252))

            y = 200
            for line in body_lines[:12]:
                draw.text((100, y), f"-  {line}", fill=(203, 213, 225))
                y += 65

            out_path = os.path.join(out_dir, f"slide_{i}.png")
            img.save(out_path, "PNG")
            image_paths.append(out_path)
        return image_paths
    except Exception as e:
        print(f"[RENDERER] PIL fallback error: {e}")
        return []


def render_deck_slides(file_path):
    """
    Renders any presentation deck (.pptx, .pdf) into slide PNGs and returns a list of filepaths.
    """
    ext = os.path.splitext(file_path)[1].lower()
    out_dir = get_deck_cache_dir(file_path)

    if ext == ".pdf":
        return render_pdf_to_images(file_path, out_dir)
    elif ext == ".pptx":
        return render_pptx_to_images(file_path, out_dir)
    return []


def get_slide_image_path(file_path, slide_number):
    """
    Returns the absolute path to the rendered slide image for a given slide number.
    Returns None if image does not exist.
    """
    out_dir = get_deck_cache_dir(file_path)
    img_path = os.path.join(out_dir, f"slide_{slide_number}.png")
    if os.path.exists(img_path):
        return img_path
    return None


def clear_deck_cache(file_path=None):
    """Clears cached slide images for a specific file or all files."""
    if file_path:
        out_dir = get_deck_cache_dir(file_path)
        if os.path.exists(out_dir):
            shutil.rmtree(out_dir, ignore_errors=True)
    else:
        if os.path.exists(CACHE_DIR):
            shutil.rmtree(CACHE_DIR, ignore_errors=True)
            os.makedirs(CACHE_DIR, exist_ok=True)
